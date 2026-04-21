"""
agent.py — Quadrogent core agent.

Root causes fixed in this version:
  1. write_file in work mode causes files to land in uploads/ instead of /workspace/.
     FIXED: write_file removed from WORK_TOOLS. Model must use execute_command+heredoc.

  2. Silent "stop" — model finishes streaming with no text and no tool_calls.
     LM Studio considers it done; agent loop exited. Task incomplete.
     FIXED: stream now yields {"type":"finish","reason":...}. Agent detects
     finish_reason="stop" + no tools + work mode → injects re-prompt and continues.

  3. Error context: model repeats failing commands without fixing them.
     FIXED: _make_next_step() injects exit code + stderr snippet after every tool call.

  4. Context explosion from write_file content: large file contents bloat context,
     confuse the model about what state the workspace is in.
     FIXED: write_file excluded from work mode entirely.

  5. Model loses track of workspace state mid-task.
     FIXED: Periodic workspace snapshot injected every 8 tool calls.
"""
import json
import os
import re
import shutil
from typing import Callable

from src.core.llm_client import LLMClient
from src.core.docker_manager import DockerManager
from src.core.web_search import WebSearch
from src.utils.file_manager import FileManager
from src.db.database import Database

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}


def _encode_image_message(text: str, image_path: str) -> dict:
    """Build a vision-capable user message with base64 image."""
    import base64, mimetypes, os
    ext = os.path.splitext(image_path)[1].lower()
    mime = {
        ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".png": "image/png", ".gif": "image/gif",
        ".webp": "image/webp", ".bmp": "image/bmp",
    }.get(ext, "image/jpeg")
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    content_parts = []
    if text:
        content_parts.append({"type": "text", "text": text})
    content_parts.append({
        "type": "image_url",
        "image_url": {"url": f"data:{mime};base64,{b64}"}
    })
    return {"role": "user", "content": content_parts}


# ─────────────────────────────────────────────────────────────────────────────
#  System prompts
# ─────────────────────────────────────────────────────────────────────────────

SYSTEM_WORK = """\
You are Quadrogent — an autonomous execution agent. Respond in the user's language.

ENVIRONMENT: Docker Ubuntu 22.04. Root access, internet available.
Pre-installed: python3, pip3, python3-venv, zip, unzip, curl, wget, git, jq.
Python libraries: python-pptx, python-docx, reportlab, Pillow (for creating presentations, Word docs, PDFs, images).
Project workspace: /workspace/
All files you create will be visible in the user's file explorer.

━━━ DOCUMENT GENERATION LIBRARIES ━━━
Available Python libraries for creating documents:

1. PRESENTATIONS (python-pptx):
   from pptx import Presentation
   from pptx.util import Inches, Pt
   prs = Presentation()
   slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
   title = slide.shapes.title
   title.text = "Hello World"
   prs.save('/workspace/presentation.pptx')

2. WORD DOCUMENTS (python-docx):
   from docx import Document
   from docx.shared import Inches, Pt
   doc = Document()
   doc.add_heading('Document Title', 0)
   doc.add_paragraph('A plain paragraph with some text.')
   doc.save('/workspace/document.docx')

3. PDF DOCUMENTS (reportlab):
   from reportlab.lib.pagesizes import letter
   from reportlab.pdfgen import canvas
   c = canvas.Canvas('/workspace/document.pdf', pagesize=letter)
   c.drawString(100, 750, 'Hello World')
   c.save()

4. IMAGES (Pillow):
   from PIL import Image, ImageDraw, ImageFont
   img = Image.new('RGB', (800, 600), color='white')
   draw = ImageDraw.Draw(img)
   draw.text((100, 100), 'Hello World', fill='black')
   img.save('/workspace/image.png')

When user asks to create presentations, Word documents, or PDFs:
1. Create Python script using these libraries
2. Execute it with execute_command
3. Verify file was created
4. Zip and deliver via deliver_file

━━━ CRITICAL: YOU MUST COMPLETE THE TASK ━━━
Your PRIMARY GOAL is to COMPLETE the user's task FULLY, not just describe how to do it.
You MUST keep working until deliver_file has been called successfully.
NEVER stop after making a plan. NEVER stop after describing steps.
EXECUTE EVERY STEP until the task is DONE.

━━━ MANDATORY RULES ━━━
1. ALWAYS call a tool every turn. NEVER output text without a tool call.
2. ALL files go in /workspace/ — create them with execute_command + heredoc.
3. Do NOT use write_file — it does not exist in this mode. Use execute_command.
4. After each tool result: read the output, then call the next tool immediately.
5. If exit code != 0: fix the exact error shown in stderr and retry.
6. Never repeat a failing command unchanged. Always fix it first.
7. install_packages is the ONLY way to install packages — never apt-get/pip in execute_command.
8. FINISH by zipping in /workspace/, then calling deliver_file with the zip filename.
9. DO NOT STOP until deliver_file is called. Keep going even if you think you're done.
10. ⚠️ NEVER install the same package twice! If install_packages returns [OK], the package is INSTALLED.
11. ⚠️ READ tool output! If you see [OK] or "already satisfied", DO NOT repeat the command.
12. ⚠️ CRITICAL: There is NO "uploads/" directory! NEVER use "uploads/" in ANY command!
13. ⚠️ WRONG: zip -r uploads/file.zip | RIGHT: zip -r file.zip
14. ⚠️ FILE PATHS — ALWAYS use DOUBLE QUOTES around ANY path in /workspace/:
    CORRECT:   cat "/workspace/Текстовый файл.txt"
    CORRECT:   stat "/workspace/madk1d - давно.mp3"
    WRONG:     cat /workspace/Текстовый ПРОБЕЛ файл.txt
    WRONG:     cat /workspace/madk1d-davno.mp3  (missing quotes)
    This is MANDATORY for ALL filenames, especially those with spaces, Cyrillic, or special chars.
    When in doubt: ls /workspace/ | head -20  to see the exact filename, then quote it.

━━━ HOW TO CREATE FILES (heredoc) ━━━
execute_command:
  mkdir -p /workspace/myproject/app && cat > /workspace/myproject/app/views.py << 'EOF'
  from django.shortcuts import render

  def home(request):
      return render(request, 'home.html', {})
  EOF

━━━ HOW TO CREATE DJANGO/FLASK PROJECTS ━━━
DJANGO:
  execute_command: cd /workspace && django-admin startproject mysite
  # This creates /workspace/mysite/ directory with manage.py inside
  # DO NOT use target directory argument if dir already exists!
  
FLASK:
  execute_command:
    mkdir -p /workspace/myapp && cat > /workspace/myapp/app.py << 'EOF'
    from flask import Flask
    app = Flask(__name__)
    @app.route('/')
    def hello(): return "Hello World"
    if __name__ == '__main__': app.run()
    EOF

━━━ WRONG (DO NOT DO THIS) ━━━
  ✗ write_file("portfolio/views.py", "...") — write_file does not exist here
  ✗ execute_command("pip install django") — use install_packages instead
  ✗ Stopping after making a plan — EXECUTE THE PLAN, don't just describe it
  ✗ Saying "Now we can..." — NO, JUST DO IT with execute_command
  ✗ django-admin startproject name /existing/dir — use WITHOUT target if dir exists
  ✗ zip -r result.zip . — zips everything including hidden files
  ✗ Correct: cd /workspace && zip -r myproject.zip myproject/

━━━ COMPLETION SEQUENCE ━━━
Step 1: Create project in /workspace/
Step 2: Zip it IN /workspace/ (NOT in uploads/!):
  execute_command: cd /workspace && zip -r portfolio.zip portfolio/
Step 3: Deliver the zip (just the filename):
  deliver_file: portfolio.zip

⚠️ CRITICAL: The zip MUST be in /workspace/, NOT in uploads/!
⚠️ WRONG: zip -r uploads/portfolio.zip portfolio/
✓ RIGHT: cd /workspace && zip -r portfolio.zip portfolio/

Remember: Your job is NOT to explain what needs to be done.
Your job is to DO IT by calling tools until deliver_file succeeds.
"""

SYSTEM_TALK = """\
You are Quadrogent, an open-source AI agent. Respond in the user's language.
Mode: Talk — have a normal conversation. Use markdown formatting.
"""

SYSTEM_CALC = """You are Quadrogent — a scientific computation assistant. Respond in the user's language.

MODE: Calc — you can run Python code to perform calculations, solve equations, process data.
Workspace: /workspace/ — you may read/write files there.

RULES:
1. For any numerical computation, always run Python code — never guess.
2. Use execute_command to run: python3 -c "..." or save+run a script.
3. For results: show the code AND the output clearly.
4. You may install Python packages via install_packages if needed.
5. Deliver result files with deliver_file if producing output files.
6. Do NOT use the shell for anything except running Python.
"""

SYSTEM_AUTO = """\
You are Quadrogent — autonomous AI agent. Respond in the user's language.
Workspace: /workspace/ (all files created here are visible to the user)
Pre-installed: python3, pip3, zip, unzip, curl, wget, git, python3-venv.
Python libraries: python-pptx (presentations), python-docx (Word docs), reportlab (PDFs), Pillow (images).

For document/presentation tasks:
  - Create Python script using available libraries (pptx, docx, reportlab, Pillow).
  - Execute script with execute_command.
  - Verify file creation, then zip and deliver_file.

For action tasks:
  - Call tools immediately. Do NOT just describe steps.
  - Create project files in /workspace/ using execute_command with heredoc.
  - install_packages for installing packages (never apt-get/pip in execute_command).
  - Finish with: zip project → deliver_file.

For questions: just answer with markdown.
"""

MEMORY_SUMMARIZE_PROMPT = (
    "Summarize this conversation in 1-3 sentences. What was discussed? Be brief."
)

TITLE_GEN_PROMPT_DEFAULT = (
    "Generate a short chat title (max 40 characters) based on the user message. "
    "Return ONLY the title text, no quotes, no punctuation at the end."
)

_HARD_REFORCE = (
    "❌ CRITICAL ERROR: You output text without calling a tool.\n"
    "In work mode you MUST call a tool every turn. NO EXCEPTIONS.\n"
    "DO NOT write explanations. DO NOT make plans.\n"
    "Call execute_command or install_packages RIGHT NOW to continue the task.\n"
    "The user's task is NOT complete. Keep working."
)

# Injected when model does a silent stop (finish_reason=stop, no content, no tools)
_SILENT_STOP_REFORCE = (
    "❌ ERROR: You stopped without completing the task.\n"
    "The task is NOT done. deliver_file was NOT called.\n"
    "You MUST continue working by calling the next tool RIGHT NOW.\n"
    "What is the next step to complete the task? Execute it immediately with execute_command.\n"
    "DO NOT describe what needs to be done. JUST DO IT."
)

# Injected when model tries to use uploads/ directory
_UPLOADS_ERROR = (
    "❌ CRITICAL: You used 'uploads/' in your command!\n"
    "There is NO uploads/ directory! It does NOT exist!\n"
    "NEVER use 'uploads/' in any command.\n\n"
    "WRONG: zip -r uploads/file.zip project/\n"
    "RIGHT: cd /workspace && zip -r file.zip project/\n\n"
    "Fix this RIGHT NOW by calling execute_command with the CORRECT path."
)


def _make_next_step(tool_name: str, exit_code, output_snippet: str,
                    mode: str = "auto") -> str:
    """Context-aware continue prompt injected after each tool result."""
    work_mode = (mode == "work")

    if exit_code is not None and exit_code != 0:
        return (
            f"❌ [Tool '{tool_name}' FAILED — exit code {exit_code}]\n"
            f"Error output: {output_snippet[:500]}\n\n"
            "CRITICAL: This error MUST be fixed before continuing.\n"
            "1. Read the error message carefully\n"
            "2. Identify the exact problem\n"
            "3. Call execute_command with the CORRECTED command\n"
            "DO NOT move to the next step until this is fixed.\n"
            "DO NOT explain what's wrong - FIX IT."
        )

    # web_search in non-work mode: tell model to ANSWER, not keep tool-calling
    if tool_name in ("web_search", "fetch_url") and not work_mode:
        return (
            f"✓ [web_search results returned above]\n\n"
            "You now have the search results. "
            "ANSWER the user's question directly using these results. "
            "Do NOT call any more tools. "
            "Do NOT try to read files or save anything. "
            "Just write your answer now."
        )

    if not work_mode:
        # auto/talk/calc: soft nudge, not a hard demand to keep calling tools
        return (
            f"✓ [Tool '{tool_name}' done]\n"
            "Use this result to answer the user. "
            "Only call another tool if you genuinely need more information."
        )

    # work mode: hard push to keep going
    return (
        f"✓ [Tool '{tool_name}' succeeded]\n\n"
        "Good! Now call the NEXT tool immediately to continue.\n"
        "Keep working until deliver_file has been called.\n"
        "DO NOT stop. DO NOT wait. Call the next tool NOW."
    )


def _strip_think(text: str) -> str:
    """Remove think blocks from stored content (all patterns)."""
    # Standard wrapped blocks
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<thinking>.*?</thinking>", "", text, flags=re.DOTALL | re.IGNORECASE)
    # Closing tag only: everything before </think> is thinking
    import re as _re2
    m = _re2.search(r"</think(?:ing)?>", text, _re2.IGNORECASE)
    if m:
        text = text[m.end():]
    # Unclosed opening tag: everything after is thinking
    text = re.sub(r"<think(?:ing)?>.*$", "", text, flags=re.DOTALL | re.IGNORECASE)
    return text.strip()


def _extract_exit_code(result: str):
    m = re.search(r"\[exit code:\s*(-?\d+)\]", result)
    return int(m.group(1)) if m else None


def _short_output(result: str, maxlen: int = 800) -> str:
    """Trim tool output — keep the tail (errors are at the end)."""
    clean = re.sub(r"^\[.*?\]\n?", "", result).strip()
    return ("..." + clean[-maxlen:]) if len(clean) > maxlen else clean



LANG_INSTRUCTIONS = {
    "auto": "Respond in the same language the user writes in.",
    "ru":   "Always respond in Russian (на русском языке).",
    "en":   "Always respond in English.",
    "de":   "Always respond in German (auf Deutsch).",
    "fr":   "Always respond in French (en français).",
    "es":   "Always respond in Spanish (en español).",
    "zh":   "Always respond in Chinese (用中文回答).",
    "ja":   "Always respond in Japanese (日本語で答えてください).",
}

class Agent:
    def __init__(self, db: Database):
        self.db = db
        self.llm = LLMClient()
        self.docker = DockerManager()
        self.search = WebSearch(db)
        self.files = FileManager()
        self.on_message:      Callable[[str, str], None]       | None = None
        self.on_tool_call:    Callable[[str, str, str], None]  | None = None
        self.on_stream_start: Callable[[], None]               | None = None
        self.on_stream_delta: Callable[[str], None]            | None = None
        self.on_stream_end:   Callable[[], None]               | None = None
        self.on_file_ready:   Callable[[str, str], None]       | None = None
        self.on_lm_log:       Callable[[str], None]            | None = None  # NEW: LM Studio логи
        self._stop = False
        
        # Установка callback для логирования LLM
        self.llm.on_log = self._log_lm

    def stop(self):
        self._stop = True
        self.llm.abort()

    def _log_lm(self, message: str):
        """Forward LM Studio logs to UI."""
        if self.on_lm_log:
            self.on_lm_log(message)

    def _emit(self, role: str, content: str):
        if self.on_message:
            self.on_message(role, content)

    def _emit_tool(self, name: str, args: str, result: str):
        if self.on_tool_call:
            # Guard: PyQt signal expects str; non-string args/result cause
            # "'int' object has no attribute 'strip'" downstream.
            self.on_tool_call(str(name), str(args), str(result))

    def _get_system_prompt(self, mode: str, persistent: bool = False) -> str:
        custom = self.db.get_setting(f"system_prompt_{mode}", "")
        base = custom if custom.strip() else {
            "work": SYSTEM_WORK, "talk": SYSTEM_TALK,
            "auto": SYSTEM_AUTO, "calc": SYSTEM_CALC,
        }.get(mode, SYSTEM_AUTO)
        # Only inject memories for persistent chats
        if persistent:
            memories = self.db.get_memories_text()
            if memories:
                memory_block = (
                    "━━━ PERSONAL MEMORY (facts about this user — always use these) ━━━\n"
                    + memories + "\n"
                    "━━━ END OF PERSONAL MEMORY ━━━\n\n"
                )
                base = memory_block + base
        return base

    def _get_system_with_think(self, mode: str, think_mode: bool, persistent: bool = False) -> str:
        base = self._get_system_prompt(mode, persistent=persistent)
        lang = self.db.get_setting("language", "auto")
        lang_instr = LANG_INSTRUCTIONS.get(lang, LANG_INSTRUCTIONS["auto"])
        base += f"\n\n{lang_instr}"
        if think_mode:
            base += (
                "\n\nThink mode ON: wrap ALL internal reasoning in <think>...</think> tags "
                "BEFORE your final answer. Example: <think>reasoning here</think> Answer here. "
                "The user will NOT see the <think> block — only your final answer is shown."
            )
        else:
            base += "\n\nRespond directly. Do NOT use <think> blocks or any internal reasoning tags."
        return base
    def auto_memorize(self, chat_id: int) -> bool:
        """Ask the LLM if the last exchange is worth memorising. Background-safe."""
        db_messages = self.db.get_messages(chat_id)
        # Берём только последние 4 сообщения (2 пары user-assistant), по 300 символов
        recent = [m for m in db_messages if m["role"] in ("user", "assistant")][-4:]
        if not recent:
            return False
        conversation = "\n".join(
            f"{m['role'].upper()}: {m['content'][:300]}" for m in recent
        )
        eval_messages2 = [
            {
                "role": "system",
                "content": (
                    "You are a strict personal-memory filter. "
                    "Read the conversation and decide if it contains a CONCRETE PERSONAL FACT "
                    "about the USER that is worth remembering long-term.\n\n"
                    "SAVE (save=true) ONLY if the user revealed:\n"
                    "  - Their name, age, location, occupation, language\n"
                    "  - A pet, family member, or relationship detail\n"
                    "  - A specific preference, hobby, or interest\n"
                    "  - A project name, technology stack, or ongoing task\n"
                    "  - A health detail, constraint, or lifestyle fact\n"
                    "  - Any explicit instruction like 'remember that I...'\n\n"
                    "DO NOT SAVE (save=false):\n"
                    "  - Small-talk, greetings, or rhetorical questions\n"
                    "  - Questions the user asked the AI (those are not facts about the user)\n"
                    "  - Generic statements without personal context\n"
                    "  - Duplicate facts already obvious from context\n\n"
                    "If saving, write summary as a THIRD-PERSON FACT about the user, "
                    "in the same language the user used.\n"
                    "Example: 'User\'s name is Alexey. He is a Python developer in Moscow.'\n\n"
                    'Respond ONLY with valid JSON (no markdown): '
                    '{"save": true, "summary": "..."} or {"save": false}'
                ),
            },
            {"role": "user", "content": conversation},
        ]
        try:
            memory_max_tokens = int(self.db.get_setting("memory_max_tokens", "150"))
            response = self.llm.chat(
                eval_messages2,
                use_tools=False,
                stream=False,
                max_tokens=memory_max_tokens
            )
            raw = response.get("content", "").strip()
            raw = re.sub(r"```[a-z]*|```", "", raw).strip()
            data = json.loads(raw)
            if data.get("save") and data.get("summary", "").strip():
                self.db.save_memory(chat_id, data["summary"].strip())
                return True
        except Exception:
            pass
        return False

    def _workspace_snapshot(self) -> str:
        """Get current state of /workspace for context injection."""
        exit_code, output = self.docker.execute(
            "find /workspace -not -path '*/\\.*' "
            "| head -80 | sort 2>/dev/null || echo '(empty)'"
        )
        return (
            f"[WORKSPACE STATE]\n"
            f"/workspace files:\n{output.strip()}\n\n"
            "Continue the task based on what is already created above."
        )

    # ── Tool execution ────────────────────────────────────────────────────────

    def _execute_tool(self, name: str, arguments: dict) -> str:

        if name == "execute_command":
            cmd = arguments.get("command", "")
            
            # Check for uploads/ usage
            if "uploads/" in cmd or "uploads\\" in cmd:
                error_msg = (
                    "❌ ERROR: Command contains 'uploads/' which does NOT exist!\n"
                    f"Your command: {cmd[:200]}\n\n"
                    "There is NO uploads/ directory in /workspace.\n"
                    "CORRECT: cd /workspace && zip -r myfile.zip myproject/\n"
                    "WRONG: zip -r uploads/myfile.zip myproject/"
                )
                self._emit_tool(name, cmd[:100], error_msg)
                return error_msg
            
            # Intercept raw apt-get → safe handler
            if re.search(r"(?:^|&&|\|)\s*apt(?:-get)?\s+install\b", cmd, re.MULTILINE):
                m = re.search(
                    r"apt(?:-get)?\s+install\s+(?:-[^\s]+\s+)*(.+?)(?:\s*&&|\s*\||\s*2>&1|$)",
                    cmd, re.DOTALL,
                )
                if m:
                    pkgs = [p for p in m.group(1).strip().split()
                            if not p.startswith("-") and p != "2>&1"]
                    if pkgs:
                        ec, out = self.docker.execute_apt(" ".join(pkgs))
                        result = f"[intercepted apt | exit code: {ec}]\n{out}"
                        self._emit_tool(name, cmd, result)
                        return result
            # Timestamp before command so we can detect newly created/modified files
            self.docker.execute("touch /tmp/.cmd_ts 2>/dev/null")
            ec, out = self.docker.execute(cmd)
            result = f"[exit code: {ec}]\n{out}" if out.strip() else f"[exit code: {ec}]"

            # After a successful command, show what changed in /workspace/
            # This prevents the model from recreating files it already created
            if ec == 0:
                _, new_files = self.docker.execute(
                    "find /workspace -newer /tmp/.cmd_ts "
                    "-not -path '*/.*' -not -name '.gitkeep' 2>/dev/null | sort"
                )
                if new_files.strip():
                    result += f"\n[Created/modified:\n{new_files.strip()}]"

            self._emit_tool(name, cmd, result)
            return result

        elif name == "web_search":
            query = arguments.get("query", "")
            results = self.search.search(query)
            formatted = self.search.format_results(results)
            self._emit_tool(name, query, formatted)
            return formatted

        elif name == "fetch_url":
            url = arguments.get("url", "").strip()
            if not url:
                result = "Error: no URL provided"
                self._emit_tool(name, url, result)
                return result
            try:
                from src.core.web_search import _fetch_page
                proxy = self.db.get_setting("search_proxy", "") if self.db else ""
                text = _fetch_page(url, proxy=proxy, max_chars=4000)
                self._emit_tool(name, url, f"[{len(text)} chars]")
                return text
            except Exception as e:
                result = f"Error fetching URL: {e}"
                self._emit_tool(name, url, result)
                return result

        elif name == "read_file":
            path = arguments.get("path", "")
            try:
                # Absolute /workspace/... paths live inside Docker, not on host
                if path.startswith("/workspace/"):
                    ec, out = self.docker.execute(f'cat "{path}" 2>&1')
                    if ec != 0:
                        self._emit_tool(name, path, f"Error: {out.strip()}")
                        return f"Error: {out.strip()}"
                    self._emit_tool(name, path, f"[{len(out)} chars]")
                    return out
                content = self.files.read(path)
                self._emit_tool(name, path, f"[{len(content)} chars]")
                return content
            except Exception as e:
                self._emit_tool(name, path, f"Error: {e}")
                return f"Error: {e}"

        elif name == "write_file":
            # In work mode this tool is not in the tool list, but model might
            # hallucinate it. Intercept gracefully.
            path = arguments.get("path", "").strip()
            content = arguments.get("content", "")
            if not path:
                result = (
                    "Error: write_file is not available in work mode. "
                    "Use execute_command with heredoc to create files in /workspace/."
                )
                self._emit_tool(name, path, result)
                return result
            try:
                abs_path = self.files.write(path, content)
                self._emit_tool(name, path, f"Written → {abs_path}")
                if self.on_file_ready:
                    self.on_file_ready(path, abs_path)
                return f"Written: {path} ({len(content.encode())} bytes)"
            except Exception as e:
                self._emit_tool(name, path, f"Error: {e}")
                return f"Error: {e}"

        elif name == "delete_file":
            path = arguments.get("path", "").strip()
            try:
                if path in (".", "", "/"):
                    base = os.path.abspath(self.files.base_dir)
                    deleted = []
                    for entry in os.listdir(base):
                        ep = os.path.join(base, entry)
                        shutil.rmtree(ep) if os.path.isdir(ep) else os.remove(ep)
                        deleted.append(entry)
                    result = f"Cleared: {len(deleted)} items"
                else:
                    self.files.delete(path)
                    result = f"Deleted: {path}"
                self._emit_tool(name, path, result)
                return result
            except Exception as e:
                self._emit_tool(name, path, f"Error: {e}")
                return f"Error: {e}"

        elif name == "list_files":
            try:
                ec, out = self.docker.execute(
                    "find /workspace -not -path '*/.*' -not -name '.gitkeep' "
                    "| sort | head -200 2>/dev/null"
                )
                if ec == 0 and out.strip():
                    result = f"Docker /workspace contents:\n{out.strip()}"
                else:
                    result = "workspace/ is empty"
                self._emit_tool(name, "", result)
                return result
            except Exception:
                return "workspace/ is empty"

        elif name == "install_packages":
            packages = arguments.get("packages", "").strip()
            manager  = arguments.get("manager", "apt").lower()
            if not packages:
                return "Error: no packages specified"
            if manager == "pip":
                ec, out = self.docker.execute(f"pip3 install --quiet {packages} 2>&1")
            else:
                ec, out = self.docker.execute_apt(packages)
            status = "OK" if ec == 0 else f"exit code {ec}"
            result = f"[{status}]\n{out}" if out else f"[{status}]"
            self._emit_tool(name, f"{manager}: {packages}", result)
            return result

        elif name == "deliver_file":
            path = arguments.get("path", "").strip()
            filename = os.path.basename(path) or path
            workspace_abs = os.path.abspath("workspace")
            host_dest = os.path.join(workspace_abs, filename)

            # Always pull from Docker via get_archive — works even if bind-mount is broken
            docker_path = f"/workspace/{filename}"
            ec, _ = self.docker.execute(f"test -f {docker_path}")
            if ec == 0:
                os.makedirs(workspace_abs, exist_ok=True)
                self.docker.copy_from_container(docker_path, host_dest)

            try:
                if not os.path.exists(host_dest):
                    result = f"Error: '{filename}' not found in Docker /workspace/"
                    self._emit_tool(name, filename, result)
                    return result
                self._emit_tool(name, filename, f"Delivered: {host_dest}")
                if self.on_file_ready:
                    self.on_file_ready(filename, host_dest)
                return f"DELIVERED: {filename}"
            except Exception as e:
                self._emit_tool(name, filename, f"Error: {e}")
                return f"Error: {e}"

        return f"Unknown tool: {name}"

    # ── Main agent loop ───────────────────────────────────────────────────────

    def run(self, chat_id: int, user_message: str, mode: str = "auto", web_search: bool = True, think_mode: bool = True, persistent: bool = False):
        self._stop = False
        self.db.add_message(chat_id, "user", user_message)

        # Detect attached files (supports multiple: [Файл: f1][Файл: f2]...)
        _image_path = None     # first image for vision encoding (if vision model)
        _file_injections = []  # notes about each file
        import os as _os, re as _re
        _ws = _os.path.abspath("workspace")
        _file_matches = _re.findall(r"\[Файл: (.+?)\]", user_message)

        vision_ids = getattr(self, "_vision_model_ids", set())
        cur_model  = self.llm.model or ""
        model_has_vision = (not vision_ids) or (cur_model in vision_ids)

        for _fname in _file_matches:
            _p   = _os.path.join(_ws, _fname)
            _ext = _os.path.splitext(_fname)[1].lower()
            is_image = _ext in IMAGE_EXTENSIONS

            if is_image:
                if model_has_vision and _image_path is None and _os.path.exists(_p):
                    # Vision model → will encode first image as base64 below
                    _image_path = _p
                    _file_injections.append(
                        f"[Прикреплено изображение: /workspace/{_fname}] "
                        f"Модель видит его содержимое выше."
                    )
                else:
                    # Non-vision model → file is in workspace, use it as-is
                    _file_injections.append(
                        f"[SYSTEM NOTE] Image file at /workspace/{_fname} "
                        f"(model cannot see its contents). "
                        f"Use it in the project as-is — copy or reference it."
                    )
            else:
                # Text/binary file → try to read and inject content
                _file_content = ""
                try:
                    with open(_p, "r", encoding="utf-8", errors="replace") as _f:
                        _file_content = _f.read()
                except Exception:
                    pass

                if _file_content:
                    _file_injections.append(
                        f"[Файл '{_fname}' предзагружен]\n"
                        f"Путь: /workspace/{_fname}\n"
                        f"--- СОДЕРЖИМОЕ ---\n{_file_content[:8000]}\n--- КОНЕЦ ---\n"
                        f"Не вызывай read_file — данные уже есть выше."
                    )
                else:
                    _file_injections.append(
                        f"[SYSTEM NOTE] File at /workspace/{_fname} — "
                        f"ALWAYS quote the path: cat \"/workspace/{_fname}\""
                    )

        _file_injection = "\n\n".join(_file_injections) if _file_injections else None

        system = self._get_system_with_think(mode, think_mode, persistent=persistent)
        db_messages = self.db.get_messages(chat_id)

        messages = [{"role": "system", "content": system}]
        for m in db_messages:
            if m["role"] in ("user", "assistant"):
                msg = {"role": m["role"], "content": m["content"]}
                # For the LAST user message, upgrade to vision if image attached
                if m["role"] == "user" and m["content"] == user_message and _image_path:
                    # Strip all [Файл: ...] prefixes to get clean user text
                    text_part = _re.sub(r"\[Файл: [^\]]+\]", "", m["content"]).strip()
                    msg = _encode_image_message(text_part or "Посмотри на изображение и выполни задачу.", _image_path)
                messages.append(msg)
            elif m["role"] == "tool":
                messages.append({"role": "user", "content": m["content"]})

        # Inject file path note right after user's first message (before model responds)
        if _file_injection:
            messages.append({"role": "user", "content": _file_injection})

        use_tools  = mode in ("work", "auto", "calc")
        work_mode  = mode == "work"
        self.llm._calc_mode = (mode == "calc")
        # "required" in work mode: model MUST call a tool every turn
        tool_choice = "required" if work_mode else "auto"
        delivered   = False
        tool_call_count = 0
        max_iterations  = 60
        silent_stops    = 0  # consecutive turns with no content and no tools
        _recent_calls: list[tuple] = []  # (tool_name, arg_hash) for loop detection
        _LOOP_THRESHOLD = 3  # same call 3 times = stuck

        # Inject initial workspace snapshot so model knows what already exists.
        # Prevents repeated "startproject" / "mkdir" when workspace is not empty.
        if work_mode:
            try:
                snapshot = self._workspace_snapshot()
                messages.append({"role": "user", "content": snapshot})
            except Exception:
                pass

        for iteration in range(max_iterations):
            if self._stop:
                break

            try:
                if self.on_stream_start:
                    self.on_stream_start()

                raw_content  = ""
                tool_calls   = None
                finish_reason = "stop"

                for item in self.llm.chat(
                    messages,
                    use_tools=use_tools,
                    stream=True,
                    tool_choice=tool_choice,
                    work_mode=work_mode,
                    use_web_search=web_search,
                ):
                    if self._stop:
                        break
                    t = item.get("type")
                    if t == "delta":
                        text = item.get("content", "")
                        if text:
                            raw_content += text
                            if self.on_stream_delta:
                                self.on_stream_delta(text)
                    elif t == "tool_calls":
                        tool_calls = item["tool_calls"]
                    elif t == "finish":
                        finish_reason = item.get("reason", "stop")
                        stream_error  = item.get("stream_error")

                if self.on_stream_end:
                    self.on_stream_end()

                display_content = _strip_think(raw_content)

                # ── Stream error retry ────────────────────────────────────────
                # Connection to LM Studio dropped mid-stream. The response is
                # incomplete — retry the last request once automatically.
                if stream_error and not tool_calls:
                    stream_retries = getattr(self, "_stream_retries", 0)
                    if stream_retries < 2:
                        self._stream_retries = stream_retries + 1
                        continue   # retry same messages
                    self._stream_retries = 0
                else:
                    self._stream_retries = 0

                # ── Post-delivery wrap-up ─────────────────────────────────────
                if delivered:
                    if raw_content.strip():
                        # Save RAW content (with <think> tags) so it can be restored on load
                        self.db.add_message(chat_id, "assistant", raw_content)
                    break

                # ── Silent stop detection ─────────────────────────────────────
                # Model emitted nothing and no tool calls. This happens when
                # LM Studio considers the generation "complete" spuriously.
                if not display_content.strip() and not tool_calls:
                    silent_stops += 1
                    if silent_stops >= 3:
                        break  # genuinely stuck, give up
                    msg = {"role": "user", "content": _SILENT_STOP_REFORCE}
                    messages.append(msg)
                    self.db.add_message(chat_id, "tool", _SILENT_STOP_REFORCE)
                    continue
                silent_stops = 0

                # ── Text-only response in work mode ───────────────────────────
                if use_tools and not tool_calls:
                    if raw_content.strip():
                        self.db.add_message(chat_id, "assistant", raw_content)
                        # For context: still use display_content (without <think>)
                        messages.append({"role": "assistant", "content": display_content})
                    if work_mode:
                        messages.append({"role": "user", "content": _HARD_REFORCE})
                        self.db.add_message(chat_id, "tool", _HARD_REFORCE)
                        continue
                    else:
                        break

                # Save reasoning text
                if raw_content.strip():
                    self.db.add_message(chat_id, "assistant", raw_content)
                    # For context: still use display_content (without <think>)
                    messages.append({"role": "assistant", "content": display_content})

                # ── finish_reason=length: tool call arguments got cut off ─────
                # Model hit max_tokens in the middle of generating tool args.
                # The call is malformed/empty — skip it and re-prompt to continue.
                if finish_reason == "length":
                    _LENGTH_MSG = (
                        "⚠️ Max tokens reached — your last tool call was CUT OFF and NOT executed.\n\n"
                        "CRITICAL: Do NOT repeat the same large call. Instead:\n"
                        "1. If writing a large file: use heredoc to write it in CHUNKS\n"
                        "   (e.g. first write the HTML head, then body section by section)\n"
                        "2. If reading: you may already have the data from earlier — check context.\n"
                        "3. Never re-read a file you already read earlier in this session.\n"
                        "4. Call the NEXT logical step, not the same huge step again.\n"
                        "Continue the task — do NOT start over."
                    )
                    messages.append({"role": "user", "content": _LENGTH_MSG})
                    self.db.add_message(chat_id, "tool", _LENGTH_MSG)
                    continue

                # ── Execute tool calls ────────────────────────────────────────
                for tc in tool_calls:
                    if self._stop:
                        break
                    func = tc.get("function", {})
                    name = func.get("name", "")
                    try:
                        args = json.loads(func.get("arguments", "{}"))
                    except json.JSONDecodeError:
                        args = {}

                    # Loop detection: same tool + same args repeated too many times
                    call_sig = (name, repr(sorted(args.items()) if isinstance(args, dict) else args))
                    _recent_calls.append(call_sig)
                    if len(_recent_calls) > 10:
                        _recent_calls.pop(0)
                    repeated = _recent_calls.count(call_sig)
                    if repeated >= _LOOP_THRESHOLD:
                        loop_msg = (
                            f"⛔ LOOP DETECTED: You called {name}({args}) {repeated} times in a row.\n"
                            f"This call returns the same result every time. STOP calling it.\n"
                            f"You already have the information. Use it to proceed with the task.\n"
                            f"If you need to write a large file, break it into smaller pieces using heredoc.\n"
                            f"Call the NEXT step, not the same step again."
                        )
                        messages.append({"role": "user", "content": loop_msg})
                        self.db.add_message(chat_id, "tool", loop_msg)
                        _recent_calls.clear()  # reset after warning
                        continue  # skip executing this call

                    result = self._execute_tool(name, args)
                    tool_call_count += 1

                    if name == "deliver_file" and result.startswith("DELIVERED:"):
                        delivered = True
                        tool_choice = "auto"  # allow text-only wrap-up after delivery

                    tool_msg = f"[Tool: {name}]\n{result}"
                    messages.append({"role": "user", "content": tool_msg})
                    self.db.add_message(chat_id, "tool", tool_msg, tool=name)

                    if not delivered:
                        # Inject context-aware next-step prompt after EVERY tool
                        ec = _extract_exit_code(result)
                        snippet = _short_output(result)
                        next_msg = _make_next_step(name, ec, snippet, mode=mode)
                        messages.append({"role": "user", "content": next_msg})

                        # Periodic workspace snapshot to keep model oriented
                        if tool_call_count % 4 == 0:
                            snapshot = self._workspace_snapshot()
                            messages.append({"role": "user", "content": snapshot})

                if delivered:
                    continue  # one more turn for wrap-up

            except Exception as e:
                if self.on_stream_end:
                    self.on_stream_end()
                err_str = str(e)
                # Friendly messages for common errors
                # Strip verbose HTTP URL from error message
                import re as _re
                clean_err = err_str.split(' for url:')[0].strip()
                if "Connection refused" in err_str or "ConnectionError" in err_str:
                    hint = "Провайдер недоступен. Проверьте, что сервис запущен и URL правильный."
                elif "401" in err_str or "Unauthorized" in err_str:
                    hint = "Неверный API-ключ. Проверьте ключ в настройках → Подключение."
                elif "403" in err_str or "Forbidden" in err_str:
                    hint = "Доступ запрещён. Проверьте API-ключ и права доступа."
                elif "400" in err_str or "Bad Request" in err_str:
                    hint = "Неверный запрос (400). Возможно, модель не поддерживает изображения или инструменты. Попробуйте другую модель."
                elif "429" in err_str or "rate limit" in err_str.lower():
                    hint = "Превышен лимит запросов. Подождите или смените провайдера."
                elif "timeout" in err_str.lower():
                    hint = "Время ожидания истекло. Провайдер не отвечает."
                elif "model" in err_str.lower() and ("not found" in err_str.lower() or "does not exist" in err_str.lower()):
                    hint = "Модель не найдена у провайдера. Выберите другую модель."
                else:
                    hint = ""
                error_msg = f"⚠ {clean_err}" + (f"\n\n💡 {hint}" if hint else "")
                self._emit("error", error_msg)
                self.db.add_message(chat_id, "assistant", error_msg)
                return

    def generate_title(self, user_message: str) -> str:
        """Generate a short chat title. Uses custom prompt from DB if set."""
        custom_prompt = self.db.get_setting("title_gen_prompt", "").strip()
        sys_prompt = custom_prompt or TITLE_GEN_PROMPT_DEFAULT
        messages = [
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": user_message[:500]},
        ]
        try:
            title_max_tokens = int(self.db.get_setting("title_max_tokens", "30"))
            response = self.llm.chat(
                messages,
                use_tools=False,
                stream=False,
                max_tokens=title_max_tokens
            )
            title = response.get("content", "").strip()
            # Strip quotes and leading/trailing punctuation
            title = title.strip(chr(34) + chr(39)).strip()
            return title[:40] if title else ""
        except Exception:
            return ""

    def summarize_chat(self, chat_id: int) -> str:
        db_messages = self.db.get_messages(chat_id)
        if not db_messages:
            return ""
        conversation = "\n".join(
            f"{m['role']}: {m['content'][:500]}" for m in db_messages[:30]
        )
        messages = [
            {"role": "system",  "content": MEMORY_SUMMARIZE_PROMPT},
            {"role": "user",    "content": conversation},
        ]
        try:
            response = self.llm.chat(
                messages,
                use_tools=False,
                stream=False,
                max_tokens=300  # summary чата — пара абзацев
            )
            summary = response.get("content", "").strip()
            if summary:
                self.db.save_memory(chat_id, summary)
            return summary
        except Exception:
            return ""
